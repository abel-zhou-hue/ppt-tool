from __future__ import annotations

import asyncio
import io
import logging

import httpx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.schemas.deck import Deck

logger = logging.getLogger(__name__)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
BLANK_LAYOUT_INDEX = 6
FALLBACK_BG = RGBColor(0xF5, 0xF5, 0xF5)
FALLBACK_TEXT = RGBColor(0x33, 0x33, 0x33)


async def assemble_pptx(deck: Deck) -> bytes:
    image_cache = await _fetch_all_images(deck)
    return _render_pptx(deck, image_cache)


async def _fetch_all_images(deck: Deck) -> dict[str, bytes]:
    urls = {s.image_url for s in deck.slides if s.image_url}
    if not urls:
        return {}

    async with httpx.AsyncClient(timeout=30.0) as client:
        async def _fetch(url: str):
            try:
                r = await client.get(url)
                r.raise_for_status()
                return url, r.content
            except Exception:
                logger.exception(f"Failed to fetch image: {url}")
                return url, None

        results = await asyncio.gather(*[_fetch(u) for u in urls])
    return {url: data for url, data in results if data is not None}


def _render_pptx(deck: Deck, image_cache: dict[str, bytes]) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[BLANK_LAYOUT_INDEX]

    for slide_data in deck.slides:
        slide = prs.slides.add_slide(blank)
        image_bytes = (
            image_cache.get(slide_data.image_url) if slide_data.image_url else None
        )
        if image_bytes:
            slide.shapes.add_picture(
                io.BytesIO(image_bytes),
                0,
                0,
                width=SLIDE_WIDTH,
                height=SLIDE_HEIGHT,
            )
        else:
            _render_fallback(slide, slide_data.slide_script)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _render_fallback(slide, slide_script: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = FALLBACK_BG

    tb = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(11.333), Inches(2.5)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"[图像未生成]\n\n{slide_script}"
    run.font.size = Pt(20)
    run.font.color.rgb = FALLBACK_TEXT
